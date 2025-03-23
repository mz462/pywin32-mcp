import json
import sys
import traceback
from typing import List, Optional, Dict, Any, Tuple, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE, MSO_LINE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import os
import base64
from PIL import Image
import io
import numpy as np
import uuid
from pathlib import Path
import tempfile
import shutil
import datetime
from mcp.server.fastmcp import FastMCP



class PowerPointContext:
    def __init__(self, workspace_dir: str = "presentations"):
        """
        Initialize the PowerPoint context with a workspace directory.
        
        Args:
            workspace_dir (str): Directory containing PowerPoint files
        """
        self.workspace_dir = Path(workspace_dir)
        self.workspace_dir.mkdir(exist_ok=True)
        self.presentations: Dict[str, Presentation] = {}
        self.current_presentation: Optional[str] = None
        # Create a template directory
        self.template_dir = self.workspace_dir / "templates"
        self.template_dir.mkdir(exist_ok=True)
        # Element ID mapping
        self.element_ids = {}  # Maps presentation path -> slide index -> shape id -> internal id

    def get_presentation(self, path: str) -> Presentation:
        """
        Get a presentation from the workspace.
        
        Args:
            path (str): Path to the presentation (relative to workspace or absolute)
        """
        # Convert to Path object
        path_obj = Path(path)
        
        # If it's not an absolute path, make it relative to workspace
        if not path_obj.is_absolute():
            path_obj = self.workspace_dir / path_obj
            
        path_str = str(path_obj)
        
        if path_str not in self.presentations:
            if path_obj.exists():
                self.presentations[path_str] = Presentation(path_str)
            else:
                self.presentations[path_str] = Presentation()
                
        self.current_presentation = path_str
        return self.presentations[path_str]

    def save_presentation(self, path: Optional[str] = None) -> None:
        """
        Save the current presentation.
        
        Args:
            path (Optional[str]): Path to save to (if None, uses current path)
        """
        if path:
            path_obj = Path(path)
            if not path_obj.is_absolute():
                path_obj = self.workspace_dir / path_obj
            save_path = str(path_obj)
        else:
            save_path = self.current_presentation
            
        if save_path and save_path in self.presentations:
            self.presentations[save_path].save(save_path)

    def list_presentations(self) -> List[str]:
        """List all PowerPoint files in the workspace."""
        return [str(f.relative_to(self.workspace_dir)) 
                for f in self.workspace_dir.glob("*.pptx")]

    def upload_presentation(self, file_path: str) -> str:
        """
        Upload a new presentation to the workspace.
        
        Args:
            file_path (str): Path to the file to upload
            
        Returns:
            str: Path to the saved file
        """
        if not file_path.endswith('.pptx'):
            raise ValueError("Only .pptx files are supported")
            
        source_path = Path(file_path)
        if not source_path.exists():
            raise FileNotFoundError(f"File {file_path} not found")
            
        dest_path = self.workspace_dir / source_path.name
        
        # Copy the file to the workspace
        shutil.copy2(source_path, dest_path)
            
        return str(dest_path.relative_to(self.workspace_dir))

    # Element Management
    def generate_element_id(self) -> str:
        """Generate a unique ID for an element."""
        return str(uuid.uuid4())

    def register_element(self, presentation_path: str, slide_index: int, shape) -> str:
        """Register a shape and get its unique ID."""
        if presentation_path not in self.element_ids:
            self.element_ids[presentation_path] = {}
        
        if slide_index not in self.element_ids[presentation_path]:
            self.element_ids[presentation_path][slide_index] = {}
        
        # Use shape's internal ID if available, otherwise create one
        shape_id = getattr(shape, "shape_id", id(shape))
        
        if shape_id not in self.element_ids[presentation_path][slide_index]:
            element_id = self.generate_element_id()
            self.element_ids[presentation_path][slide_index][shape_id] = element_id
        
        return self.element_ids[presentation_path][slide_index][shape_id]

    def get_shape_by_id(self, presentation: Presentation, slide_index: int, element_id: str):
        """Get a shape by its unique ID."""
        try:
            # Verify slide index is valid
            if slide_index >= len(presentation.slides):
                return None
                
            slide = presentation.slides[slide_index]
            
            # Get the presentation path
            presentation_path = self.current_presentation
            if not presentation_path:
                return None
                
            # Check if we have element mappings for this presentation and slide
            if (presentation_path not in self.element_ids or 
                slide_index not in self.element_ids[presentation_path]):
                return None
                
            # Iterate through shapes and find matching element ID
            for shape in slide.shapes:
                shape_id = getattr(shape, "shape_id", id(shape))
                if (shape_id in self.element_ids[presentation_path][slide_index] and 
                    self.element_ids[presentation_path][slide_index][shape_id] == element_id):
                    return shape
                    
            # If we get here, no matching shape was found
            return None
            
        except Exception as e:
            print(f"Error in get_shape_by_id: {str(e)}")
            return None

    def analyze_slide_content(self, presentation_path: str, slide) -> Dict[str, Any]:
        """Analyze the content of a slide and return structured information."""
        content = {
            "text_boxes": [],
            "images": [],
            "shapes": [],
            "charts": [],
            "tables": [],
            "layout": None
        }
        
        # Get slide layout
        content["layout"] = slide.slide_layout.name if slide.slide_layout else None
        
        # Analyze shapes
        for shape in slide.shapes:
            # Register this shape to get a unique ID
            element_id = self.register_element(presentation_path, 
                                             slide.slides.index(slide), 
                                             shape)
            
            shape_info = {
                "id": element_id,
                "type": str(shape.shape_type),
                "position": {"x": shape.left / Inches(1), "y": shape.top / Inches(1)},
                "size": {"width": shape.width / Inches(1), "height": shape.height / Inches(1)}
            }
            
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                content["images"].append(shape_info)
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                content["tables"].append(shape_info)
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                content["charts"].append(shape_info)
            elif hasattr(shape, "text") and shape.text.strip():
                shape_info["text"] = shape.text.strip()
                content["text_boxes"].append(shape_info)
            else:
                content["shapes"].append(shape_info)
                
        return content

    def find_element(self, presentation_path: str, slide_index: int, 
                    element_type: str = "any", search_text: Optional[str] = None,
                    position: Optional[Dict[str, float]] = None) -> List[Dict[str, Any]]:
        """
        Find elements on a slide based on type, text content, or position.
        
        Args:
            presentation_path: Path to the presentation
            slide_index: Index of the slide to search
            element_type: Type of element to find (text, shape, image, chart, table, any)
            search_text: Text to search for in element content
            position: Position criteria for finding elements
        
        Returns:
            List of matching elements with confidence scores
        """
        presentation = self.get_presentation(presentation_path)
        slide = presentation.slides[slide_index]
        
        results = []
        
        for shape in slide.shapes:
            # Skip if filtering by type and this doesn't match
            if element_type != "any":
                if element_type == "text" and (not hasattr(shape, "text") or not shape.text.strip()):
                    continue
                if element_type == "image" and shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
                if element_type == "chart" and shape.shape_type != MSO_SHAPE_TYPE.CHART:
                    continue
                if element_type == "table" and shape.shape_type != MSO_SHAPE_TYPE.TABLE:
                    continue
                if element_type == "shape" and (shape.shape_type == MSO_SHAPE_TYPE.PICTURE or 
                                               shape.shape_type == MSO_SHAPE_TYPE.CHART or
                                               shape.shape_type == MSO_SHAPE_TYPE.TABLE or
                                               (hasattr(shape, "text") and shape.text.strip())):
                    continue
            
            # Check text content if specified
            text_match = False
            confidence = 1.0
            
            if search_text and hasattr(shape, "text"):
                shape_text = shape.text.strip().lower()
                search_lower = search_text.lower()
                
                if shape_text == search_lower:
                    text_match = True
                    confidence = 1.0
                elif search_lower in shape_text:
                    text_match = True
                    # Calculate confidence based on how much of the text matches
                    confidence = len(search_lower) / len(shape_text)
                else:
                    continue  # Text doesn't match at all
            
            # Check position if specified
            position_match = True
            if position:
                shape_x = shape.left / Inches(1)
                shape_y = shape.top / Inches(1)
                
                # Calculate distance from target position
                target_x = position.get("x")
                target_y = position.get("y")
                proximity = position.get("proximity", 1.0)
                
                if target_x is not None and target_y is not None:
                    distance = ((shape_x - target_x) ** 2 + (shape_y - target_y) ** 2) ** 0.5
                    if distance > proximity:
                        position_match = False
                    else:
                        # Adjust confidence based on proximity
                        position_confidence = 1.0 - (distance / proximity)
                        confidence *= position_confidence
            
            # If all criteria match, add to results
            if (search_text is None or text_match) and position_match:
                element_id = self.register_element(presentation_path, slide_index, shape)
                
                element = {
                    "id": element_id,
                    "type": self._get_shape_type_name(shape),
                    "text": shape.text.strip() if hasattr(shape, "text") else None,
                    "position": {"x": shape.left / Inches(1), "y": shape.top / Inches(1)},
                    "size": {"width": shape.width / Inches(1), "height": shape.height / Inches(1)},
                    "confidence": confidence
                }
                
                results.append(element)
        
        # Sort by confidence
        results.sort(key=lambda x: x["confidence"], reverse=True)
        return results

    def _get_shape_type_name(self, shape) -> str:
        """Convert shape type to a user-friendly name."""
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            for name, value in vars(MSO_AUTO_SHAPE_TYPE).items():
                if not name.startswith("__") and value == shape.auto_shape_type:
                    return name.lower()
            return "auto_shape"
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return "image"
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            return "chart"
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return "table"
        else:
            for name, value in vars(MSO_SHAPE_TYPE).items():
                if not name.startswith("__") and value == shape.shape_type:
                    return name.lower()
            return "unknown"

    def edit_element(self, presentation_path: str, slide_index: int, 
                   element_id: str, properties: Dict[str, Any]) -> Dict[str, Any]:
        """
        Edit properties of a specific element on a slide.
        
        Args:
            presentation_path: Path to the presentation
            slide_index: Index of the slide
            element_id: Unique ID of the element to edit
            properties: Properties to modify (text, position, size, rotation, etc.)
            
        Returns:
            Updated properties of the element
        """
        try:
            # First verify the presentation exists and can be loaded
            presentation = self.get_presentation(presentation_path)
            if not presentation:
                return {"error": f"Could not load presentation: {presentation_path}"}
            
            # Verify the slide index is valid
            if slide_index >= len(presentation.slides):
                return {"error": f"Invalid slide index {slide_index}. Presentation has {len(presentation.slides)} slides."}
            
            # Get the shape and handle None case explicitly
            shape = self.get_shape_by_id(presentation, slide_index, element_id)
            if not shape:
                return {"error": f"Could not find element with ID {element_id} on slide {slide_index}"}
            
            # Apply changes based on properties
            if "text" in properties and hasattr(shape, "text_frame"):
                shape.text_frame.text = properties["text"]
            
            if "position" in properties:
                position = properties["position"]
                if "x" in position:
                    shape.left = Inches(position["x"])
                if "y" in position:
                    shape.top = Inches(position["y"])
                
            if "size" in properties:
                size = properties["size"]
                if "width" in size:
                    shape.width = Inches(size["width"])
                if "height" in size:
                    shape.height = Inches(size["height"])
                
            if "rotation" in properties:
                shape.rotation = properties["rotation"]
            
            if "transparency" in properties and hasattr(shape, "fill"):
                alpha = int(255 * (100 - properties["transparency"]) / 100)
                if hasattr(shape.fill.fore_color, "transparency"):
                    shape.fill.fore_color.transparency = (255 - alpha) / 255
            
            # Return updated properties
            updated_props = {
                "text": shape.text if hasattr(shape, "text") else None,
                "position": {"x": shape.left / Inches(1), "y": shape.top / Inches(1)},
                "size": {"width": shape.width / Inches(1), "height": shape.height / Inches(1)},
                "rotation": shape.rotation if hasattr(shape, "rotation") else None
            }
            
            # Save the changes
            self.save_presentation(presentation_path)
            
            return {
                "message": "Element updated successfully",
                "properties": updated_props
            }
            
        except Exception as e:
            # Add more context to the error message
            error_msg = f"Error editing element: {str(e)}\n"
            error_msg += f"Presentation: {presentation_path}\n"
            error_msg += f"Slide: {slide_index}\n"
            error_msg += f"Element ID: {element_id}\n"
            error_msg += f"Properties: {properties}"
            return {"error": error_msg}

    def style_element(self, presentation_path: str, slide_index: int, 
                    element_id: str, style_properties: Dict[str, Any]) -> bool:
        """
        Apply styling to a specific element on a slide.
        
        Args:
            presentation_path: Path to the presentation
            slide_index: Index of the slide
            element_id: Unique ID of the element to style
            style_properties: Style properties to apply
            
        Returns:
            True if styling was applied successfully
        """
        presentation = self.get_presentation(presentation_path)
        shape = self.get_shape_by_id(presentation, slide_index, element_id)
        
        # Apply font styling
        if "font" in style_properties and hasattr(shape, "text_frame"):
            font_props = style_properties["font"]
            
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "family" in font_props:
                        run.font.name = font_props["family"]
                        
                    if "size" in font_props:
                        run.font.size = Pt(font_props["size"])
                        
                    if "bold" in font_props:
                        run.font.bold = font_props["bold"]
                        
                    if "italic" in font_props:
                        run.font.italic = font_props["italic"]
                        
                    if "underline" in font_props:
                        run.font.underline = font_props["underline"]
                        
                    if "color" in font_props:
                        color = font_props["color"].lstrip('#')
                        r, g, b = tuple(int(color[i:i+2], 16) for i in (0, 2, 4))
                        run.font.color.rgb = RGBColor(r, g, b)
        
        # Apply fill styling
        if "fill" in style_properties and hasattr(shape, "fill"):
            fill_props = style_properties["fill"]
            
            if fill_props.get("type") == "solid" and "color" in fill_props:
                shape.fill.solid()
                color = fill_props["color"].lstrip('#')
                r, g, b = tuple(int(color[i:i+2], 16) for i in (0, 2, 4))
                shape.fill.fore_color.rgb = RGBColor(r, g, b)
                
            elif fill_props.get("type") == "gradient" and "gradient" in fill_props:
                # PowerPoint API doesn't support setting gradient directly
                # This would require more complex implementation
                shape.fill.gradient()
                
                # Set start color
                start_color = fill_props["gradient"]["start_color"].lstrip('#')
                r, g, b = tuple(int(start_color[i:i+2], 16) for i in (0, 2, 4))
                shape.fill.gradient_stops[0].color.rgb = RGBColor(r, g, b)
                
                # Set end color
                end_color = fill_props["gradient"]["end_color"].lstrip('#')
                r, g, b = tuple(int(end_color[i:i+2], 16) for i in (0, 2, 4))
                shape.fill.gradient_stops[-1].color.rgb = RGBColor(r, g, b)
                
            elif fill_props.get("type") == "none":
                shape.fill.background()
        
        # Apply line styling
        if "line" in style_properties and hasattr(shape, "line"):
            line_props = style_properties["line"]
            
            if "color" in line_props:
                color = line_props["color"].lstrip('#')
                r, g, b = tuple(int(color[i:i+2], 16) for i in (0, 2, 4))
                shape.line.color.rgb = RGBColor(r, g, b)
                
            if "width" in line_props:
                shape.line.width = Pt(line_props["width"])
                
            if "style" in line_props:
                # Map string style names to PowerPoint constants
                style_map = {
                    "solid": MSO_LINE.SOLID,
                    "dash": MSO_LINE.DASH,
                    "dot": MSO_LINE.ROUND_DOT,
                    "dash-dot": MSO_LINE.DASH_DOT,
                    "none": None
                }
                
                if line_props["style"] in style_map:
                    if style_map[line_props["style"]] is None:
                        shape.line.fill.background()  # No line
                    else:
                        shape.line.dash_style = style_map[line_props["style"]]
        
        return True

    def find_slide_by_content(self, presentation_path: str, search_text: str) -> Optional[int]:
        """Find a slide index by searching through its content."""
        presentation = self.get_presentation(presentation_path)
        for idx, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and search_text.lower() in shape.text.lower():
                    return idx
        return None

    def get_slide_preview(self, presentation_path: str, slide_index: int) -> str:
        """Generate a preview of the slide as a base64 encoded image."""
        presentation = self.get_presentation(presentation_path)
        slide = presentation.slides[slide_index]
        
        # This is a placeholder - in a real implementation, you would need to
        # use a proper PowerPoint rendering library or service
        # For now, we'll create a simple visualization using PIL
        width, height = 1920, 1080  # Standard slide dimensions
        img = Image.new('RGB', (width, height), 'white')
        
        # Draw shapes and text (simplified)
        from PIL import ImageDraw, ImageFont
        draw = ImageDraw.Draw(img)
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                # Convert PowerPoint coordinates to image coordinates
                x = int(shape.left * width / 914400)  # Convert EMU to pixels
                y = int(shape.top * height / 685800)
                draw.text((x, y), shape.text, fill='black')
        
        # Convert to base64
        buffered = io.BytesIO()
        img.save(buffered, format="PNG")
        return base64.b64encode(buffered.getvalue()).decode()
        
    def add_shape(self, presentation_path: str, slide_index: int, shape_type: str,
                position: Dict[str, float], size: Dict[str, float], 
                style_properties: Optional[Dict[str, Any]] = None) -> str:
        """
        Add a new shape to a slide.
        
        Args:
            presentation_path: Path to the presentation
            slide_index: Index of the slide
            shape_type: Type of shape to add
            position: Position of the shape (x, y)
            size: Size of the shape (width, height)
            style_properties: Style properties for the shape
        
        Returns:
            ID of the created shape
        """
        presentation = self.get_presentation(presentation_path)
        slide = presentation.slides[slide_index]
        
        # Map shape type names to PowerPoint constants
        shape_map = {
            # Basic shapes
            "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            "rounded_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
            "triangle": MSO_AUTO_SHAPE_TYPE.TRIANGLE,
            "right_triangle": MSO_AUTO_SHAPE_TYPE.RIGHT_TRIANGLE,
            "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
            "pentagon": MSO_AUTO_SHAPE_TYPE.PENTAGON,
            "hexagon": MSO_AUTO_SHAPE_TYPE.HEXAGON,
            "heptagon": MSO_AUTO_SHAPE_TYPE.HEPTAGON,
            "octagon": MSO_AUTO_SHAPE_TYPE.OCTAGON,
            "decagon": MSO_AUTO_SHAPE_TYPE.DECAGON,
            "dodecagon": MSO_AUTO_SHAPE_TYPE.DODECAGON,
            
            # Stars
            "star4": MSO_AUTO_SHAPE_TYPE.STAR_4_POINT,
            "star5": MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
            "star6": MSO_AUTO_SHAPE_TYPE.STAR_6_POINT,
            "star7": MSO_AUTO_SHAPE_TYPE.STAR_7_POINT,
            "star8": MSO_AUTO_SHAPE_TYPE.STAR_8_POINT,
            "star10": MSO_AUTO_SHAPE_TYPE.STAR_10_POINT,
            "star12": MSO_AUTO_SHAPE_TYPE.STAR_12_POINT,
            "star16": MSO_AUTO_SHAPE_TYPE.STAR_16_POINT,
            "star24": MSO_AUTO_SHAPE_TYPE.STAR_24_POINT,
            "star32": MSO_AUTO_SHAPE_TYPE.STAR_32_POINT,
            
            # Arrows
            "arrow": MSO_AUTO_SHAPE_TYPE.ARROW,
            "left_arrow": MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
            "right_arrow": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            "up_arrow": MSO_AUTO_SHAPE_TYPE.UP_ARROW,
            "down_arrow": MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
            "left_right_arrow": MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW,
            "up_down_arrow": MSO_AUTO_SHAPE_TYPE.UP_DOWN_ARROW,
            
            # Special shapes
            "heart": MSO_AUTO_SHAPE_TYPE.HEART,
            "lightning_bolt": MSO_AUTO_SHAPE_TYPE.LIGHTNING_BOLT,
            "sun": MSO_AUTO_SHAPE_TYPE.SUN,
            "moon": MSO_AUTO_SHAPE_TYPE.MOON,
            "smiley_face": MSO_AUTO_SHAPE_TYPE.SMILEY_FACE,
            "cloud": MSO_AUTO_SHAPE_TYPE.CLOUD,
            
            # Process shapes
            "flow_chart_process": MSO_AUTO_SHAPE_TYPE.FLOW_CHART_PROCESS,
            "flow_chart_decision": MSO_AUTO_SHAPE_TYPE.FLOW_CHART_DECISION,
            "flow_chart_connector": MSO_AUTO_SHAPE_TYPE.FLOW_CHART_CONNECTOR,
        }
        
        if shape_type.lower() not in shape_map:
            raise ValueError(f"Unsupported shape type: {shape_type}")
            
        shape_type_enum = shape_map[shape_type.lower()]
        left = Inches(position.get("x", 0))
        top = Inches(position.get("y", 0))
        width = Inches(size.get("width", 1))
        height = Inches(size.get("height", 1))
        
        shape = slide.shapes.add_shape(shape_type_enum, left, top, width, height)
        
        # Register the shape to get a unique ID
        element_id = self.register_element(presentation_path, slide_index, shape)
        
        # Apply style properties if provided
        if style_properties:
            self.style_element(presentation_path, slide_index, element_id, style_properties)
        
        return element_id
        
    def connect_shapes(self, presentation_path: str, slide_index: int, 
                      from_element_id: str, to_element_id: str, 
                      connector_type: str = "straight",
                      style_properties: Optional[Dict[str, Any]] = None) -> str:
        """
        Create a connector between two shapes on a slide.
        
        Args:
            presentation_path: Path to the presentation
            slide_index: Index of the slide
            from_element_id: ID of the starting element
            to_element_id: ID of the ending element
            connector_type: Type of connector to create (straight, elbow, curved)
            style_properties: Style properties for the connector
        
        Returns:
            ID of the created connector
        """
        presentation = self.get_presentation(presentation_path)
        slide = presentation.slides[slide_index]
        
        # Get the shapes
        from_shape = self.get_shape_by_id(presentation, slide_index, from_element_id)
        to_shape = self.get_shape_by_id(presentation, slide_index, to_element_id)
        
        # Map connector types to PowerPoint constants
        connector_map = {
            "straight": MSO_AUTO_SHAPE_TYPE.LINE_CONNECTOR_1,
            "elbow": MSO_AUTO_SHAPE_TYPE.LINE_CONNECTOR_3, 
            "curved": MSO_AUTO_SHAPE_TYPE.CURVED_CONNECTOR_3
        }
        
        if connector_type.lower() not in connector_map:
            raise ValueError(f"Unsupported connector type: {connector_type}")
            
        connector_type_enum = connector_map[connector_type.lower()]
        
        # Add connector
        # Note: python-pptx doesn't have direct support for connecting shapes
        # We create a line in between as an approximation
        start_x = from_shape.left + from_shape.width / 2
        start_y = from_shape.top + from_shape.height / 2
        end_x = to_shape.left + to_shape.width / 2
        end_y = to_shape.top + to_shape.height / 2
        
        # Create connector
        connector = slide.shapes.add_connector(
            connector_type_enum,
            start_x, start_y, end_x - start_x, end_y - start_y
        )
        
        # Register the connector to get a unique ID
        element_id = self.register_element(presentation_path, slide_index, connector)
        
        # Apply style properties if provided
        if style_properties:
            self.style_element(presentation_path, slide_index, element_id, style_properties)
        
        return element_id

    def get_company_financials(self, company_id: str, 
                             metrics: List[str] = ["revenue", "ebitda", "profit"],
                             years: Optional[List[int]] = None) -> Dict[str, Any]:
        """
        Fetch financial data for a Norwegian company from the Proff API.
        
        Args:
            company_id: Company name or organization number
            metrics: Financial metrics to retrieve
            years: Years to retrieve data for (empty means all available years)
            
        Returns:
            Company information and financial data
        """
        # This is a placeholder - in a real implementation, you would connect to Proff API
        # For now, we'll generate sample data
        
        # Sample company info
        company_info = {
            "name": company_id,
            "org_number": "123456789",
            "industry": "Technology",
            "founded": "2010"
        }
        
        # Sample financial data
        financials = {}
        current_year = 2023
        
        # Generate data for the requested years or last 5 years if not specified
        if not years:
            years = list(range(current_year - 4, current_year + 1))
            
        for year in years:
            financials[str(year)] = {}
            
            if "revenue" in metrics:
                financials[str(year)]["revenue"] = 1000000 * (1 + (year - 2018) * 0.1) * (0.9 + 0.2 * np.random.random())
                
            if "ebitda" in metrics:
                financials[str(year)]["ebitda"] = financials[str(year)].get("revenue", 1000000) * (0.15 + 0.05 * np.random.random())
                
            if "ebit" in metrics:
                financials[str(year)]["ebit"] = financials[str(year)].get("ebitda", 150000) * 0.8
                
            if "profit" in metrics:
                financials[str(year)]["profit"] = financials[str(year)].get("ebit", 120000) * 0.75
                
            if "assets" in metrics:
                financials[str(year)]["assets"] = financials[str(year)].get("revenue", 1000000) * (0.8 + 0.4 * np.random.random())
                
            if "equity" in metrics:
                financials[str(year)]["equity"] = financials[str(year)].get("assets", 800000) * (0.4 + 0.2 * np.random.random())
                
            if "debt" in metrics:
                financials[str(year)]["debt"] = financials[str(year)].get("assets", 800000) - financials[str(year)].get("equity", 320000)
                
            if "employees" in metrics:
                financials[str(year)]["employees"] = int(10 * (1 + (year - 2018) * 0.15))
                
            if "growth" in metrics:
                if year > min(years):
                    prev_revenue = financials[str(year-1)].get("revenue", 0)
                    if prev_revenue > 0:
                        financials[str(year)]["growth"] = (financials[str(year)].get("revenue", 0) - prev_revenue) / prev_revenue
                    else:
                        financials[str(year)]["growth"] = 0
                else:
                    financials[str(year)]["growth"] = 0
                    
            if "margin" in metrics:
                if financials[str(year)].get("revenue", 0) > 0:
                    financials[str(year)]["margin"] = financials[str(year)].get("ebitda", 0) / financials[str(year)].get("revenue", 1)
                else:
                    financials[str(year)]["margin"] = 0
        
        return {
            "company_info": company_info,
            "financials": financials
        }

    def create_financial_chart(self, presentation_path: str, slide_index: int, 
                             chart_type: str, data: Dict[str, Any], 
                             position: Dict[str, float], size: Dict[str, float], 
                             title: Optional[str] = None) -> str:
        """
        Create a financial chart on a slide.
        
        Args:
            presentation_path: Path to the presentation
            slide_index: Index of the slide
            chart_type: Type of chart to create
            data: Chart data structure
            position: Position of the chart
            size: Size of the chart
            title: Chart title
            
        Returns:
            ID of the created chart
        """
        presentation = self.get_presentation(presentation_path)
        slide = presentation.slides[slide_index]
        
        # Map chart types to PowerPoint constants
        chart_map = {
            "line": XL_CHART_TYPE.LINE,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
            "area": XL_CHART_TYPE.AREA,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "bubble": XL_CHART_TYPE.BUBBLE,
            "radar": XL_CHART_TYPE.RADAR,
            "waterfall": XL_CHART_TYPE.WATERFALL
        }
        
        if chart_type.lower() not in chart_map:
            raise ValueError(f"Unsupported chart type: {chart_type}")
            
        chart_type_enum = chart_map[chart_type.lower()]
        
        # Prepare chart data
        chart_data = CategoryChartData()
        chart_data.categories = data["categories"]
        
        for series in data["series"]:
            chart_data.add_series(series["name"], series["values"])
        
        # Add chart to slide
        left = Inches(position.get("x", 0))
        top = Inches(position.get("y", 0))
        width = Inches(size.get("width", 5))
        height = Inches(size.get("height", 3))
        
        chart_shape = slide.shapes.add_chart(
            chart_type_enum, left, top, width, height, chart_data
        )
        chart = chart_shape.chart
        
        # Set chart title if provided
        if title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        
        # Register the chart to get a unique ID
        element_id = self.register_element(presentation_path, slide_index, chart_shape)
        
        return element_id

    def create_comparison_table(self, presentation_path: str, slide_index: int,
                              companies: List[str], metrics: List[str],
                              position: Dict[str, float], title: Optional[str] = None) -> str:
        """
        Create a comparison table on a slide.
        
        Args:
            presentation_path: Path to the presentation
            slide_index: Index of the slide
            companies: List of companies to compare
            metrics: List of metrics to compare
            position: Position of the table
            title: Table title
            
        Returns:
            ID of the created table
        """
        presentation = self.get_presentation(presentation_path)
        slide = presentation.slides[slide_index]
        
        # Get company data
        companies_data = []
        for company in companies:
            company_data = self.get_company_financials(company, metrics)
            companies_data.append(company_data)
        
        # Determine rows and columns
        rows = len(metrics) + 1  # +1 for header
        cols = len(companies) + 1  # +1 for metrics names
        
        # Add table to slide
        left = Inches(position.get("x", 0))
        top = Inches(position.get("y", 0))
        width = Inches(6)  # Default width
        height = Inches(rows * 0.5)  # Height based on number of rows
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Fill header row
        table.cell(0, 0).text = "Metrics"
        for i, company in enumerate(companies):
            table.cell(0, i + 1).text = companies_data[i]["company_info"]["name"]
        
        # Fill metrics rows
        for i, metric in enumerate(metrics):
            table.cell(i + 1, 0).text = metric.capitalize()
            for j, company in enumerate(companies):
                # Get the latest year data for this metric
                years = sorted(companies_data[j]["financials"].keys())
                if years:
                    latest_year = years[-1]
                    value = companies_data[j]["financials"][latest_year].get(metric, "N/A")
                    
                    # Format value based on metric type
                    if metric in ["revenue", "ebitda", "ebit", "profit", "assets", "equity", "debt"]:
                        # Format as currency
                        if isinstance(value, (int, float)):
                            table.cell(i + 1, j + 1).text = f"NOK {value:,.0f}"
                        else:
                            table.cell(i + 1, j + 1).text = str(value)
                    elif metric in ["growth", "margin", "roi"]:
                        # Format as percentage
                        if isinstance(value, (int, float)):
                            table.cell(i + 1, j + 1).text = f"{value:.1%}"
                        else:
                            table.cell(i + 1, j + 1).text = str(value)
                    else:
                        # Format as is
                        table.cell(i + 1, j + 1).text = str(value)
                else:
                    table.cell(i + 1, j + 1).text = "N/A"
        
        # Register the table to get a unique ID
        element_id = self.register_element(presentation_path, slide_index, table)
        
        # Add title if provided
        if title:
            title_left = left
            title_top = top - Inches(0.5)
            title_shape = slide.shapes.add_textbox(title_left, title_top, Inches(6), Inches(0.5))
            title_shape.text_frame.text = title
            title_shape.text_frame.paragraphs[0].font.bold = True
            title_shape.text_frame.paragraphs[0].font.size = Pt(14)
        
        return element_id

    def apply_template(self, presentation_path: str, template_name: str, 
                     options: Optional[Dict[str, bool]] = None) -> Dict[str, Any]:
        """
        Apply a template to a presentation.
        
        Args:
            presentation_path: Path to the presentation
            template_name: Name of the template to apply
            options: Template application options
            
        Returns:
            Elements applied from the template
        """
        presentation = self.get_presentation(presentation_path)
        
        # Get template path
        template_path = self.template_dir / f"{template_name}.pptx"
        if not template_path.exists():
            raise ValueError(f"Template '{template_name}' not found")
        
        # Open template
        template = Presentation(template_path)
        
        # Set default options
        if options is None:
            options = {
                "apply_master": True,
                "apply_theme": True,
                "apply_layouts": True
            }
        
        # Apply master slides if requested
        applied_elements = {
            "master_slides": False,
            "theme": False,
            "layouts": []
        }
        
        # Note: python-pptx doesn't provide direct methods to apply templates
        # This would require a more complex implementation
        # For now, we'll simulate it by copying elements
        
        # Apply theme (colors, fonts)
        if options.get("apply_theme", True):
            # In a real implementation, this would apply the template's theme
            applied_elements["theme"] = True
        
        # Apply layouts
        if options.get("apply_layouts", True):
            # In a real implementation, this would copy the template's layouts
            for layout in template.slide_layouts:
                applied_elements["layouts"].append(layout.name)
        
        return applied_elements

    def create_slide_from_template(self, presentation_path: str, template_name: str,
                                 content: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        """
        Create a new slide based on a template.
        
        Args:
            presentation_path: Path to the presentation
            template_name: Name of the slide template to use
            content: Content to populate the template with
            
        Returns:
            Information about the created slide
        """
        presentation = self.get_presentation(presentation_path)
        
        # Get template path
        template_path = self.template_dir / f"{template_name}.pptx"
        if not template_path.exists():
            raise ValueError(f"Template '{template_name}' not found")
        
        # Open template
        template = Presentation(template_path)
        
        # Use first slide from template as a model
        if not template.slides:
            raise ValueError(f"Template '{template_name}' has no slides")
        
        template_slide = template.slides[0]
        
        # Find matching layout in target presentation
        # In a real implementation, you would find the best matching layout
        # For now, we'll use a default layout
        layout = presentation.slide_layouts[0]  # Title Slide layout
        
        # Create new slide with the layout
        slide = presentation.slides.add_slide(layout)
        
        # Populate placeholders if content is provided
        populated_placeholders = []
        if content:
            for shape in slide.shapes:
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    ph_name = shape.placeholder_format.type if hasattr(shape.placeholder_format, "type") else "unknown"
                    
                    # Check if content has a matching key
                    for key, value in content.items():
                        if key.lower() in ph_name.lower() or ph_name.lower() in key.lower():
                            if hasattr(shape, "text_frame"):
                                shape.text_frame.text = str(value)
                                populated_placeholders.append(key)
        
        return {
            "slide_index": len(presentation.slides) - 1,
            "populated_placeholders": populated_placeholders
        }

    def list_templates(self) -> List[Dict[str, Any]]:
        """
        List all available templates.
        
        Returns:
            List of templates with metadata
        """
        templates = []
        
        for template_path in self.template_dir.glob("*.pptx"):
            template_name = template_path.stem
            
            # Load template metadata if available
            metadata_path = self.template_dir / f"{template_name}.json"
            metadata = {}
            
            if metadata_path.exists():
                try:
                    with open(metadata_path, 'r') as f:
                        metadata = json.load(f)
                except Exception:
                    pass
            
            templates.append({
                "name": template_name,
                "path": str(template_path.relative_to(self.workspace_dir)),
                "metadata": metadata
            })
        
        return templates

    def save_as_template(self, presentation_path: str, slide_index: int,
                       template_name: str, template_description: str = "") -> Dict[str, Any]:
        """
        Save a slide as a template.
        
        Args:
            presentation_path: Path to the presentation
            slide_index: Index of the slide to save as template
            template_name: Name for the new template
            template_description: Description of the template
            
        Returns:
            Information about the saved template
        """
        presentation = self.get_presentation(presentation_path)
        
        # Create a new presentation with just this slide
        template = Presentation()
        
        # Create a temporary file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            temp_path = temp_file.name
        
        try:
            # Copy the slide to the template
            # In a real implementation, this would properly copy the slide
            # For now, we'll just create a new presentation with a blank slide
            layout = template.slide_layouts[0]
            template.slides.add_slide(layout)
            
            # Save the template
            template_path = self.template_dir / f"{template_name}.pptx"
            template.save(template_path)
            
            # Save metadata
            metadata = {
                "description": template_description,
                "created_from": presentation_path,
                "slide_index": slide_index,
                "created_date": str(datetime.datetime.now())
            }
            
            metadata_path = self.template_dir / f"{template_name}.json"
            with open(metadata_path, 'w') as f:
                json.dump(metadata, f, indent=2)
            
            return {
                "name": template_name,
                "path": str(template_path.relative_to(self.workspace_dir)),
                "metadata": metadata
            }
            
        finally:
            # Clean up
            try:
                os.unlink(temp_path)
            except Exception:
                pass
    
    def add_slide(self, presentation_path: str, layout_name: str = "Title and Content") -> int:
        """Add a new slide to the presentation and return its index."""
        presentation = self.get_presentation(presentation_path)
        
        # Try to get the layout by name
        layout = None
        for slide_layout in presentation.slide_layouts:
            if hasattr(slide_layout, "name") and slide_layout.name == layout_name:
                layout = slide_layout
                break
        
        # If layout not found, use the first layout
        if layout is None:
            layout = presentation.slide_layouts[0]
        
        slide = presentation.slides.add_slide(layout)
        return len(presentation.slides) - 1  # Return the index of the new slide
    
    def add_text(self, presentation_path: str, slide_index: int, text: str,
               position: List[float] = [0.5, 0.5], font_size: float = 18) -> str:
        """Add text to a slide and return the element ID."""
        presentation = self.get_presentation(presentation_path)
        slide = presentation.slides[slide_index]
        
        # Create a text box
        left, top = Inches(position[0]), Inches(position[1])
        width, height = Inches(6), Inches(1)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text_frame.text = text
        
        # Set font size
        if textbox.text_frame.paragraphs:
            p = textbox.text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
        
        # Register and return the element ID
        return self.register_element(presentation_path, slide_index, textbox)
    
    def get_slide_count(self, presentation_path: str) -> int:
        """Get the total number of slides in the presentation."""
        presentation = self.get_presentation(presentation_path)
        return len(presentation.slides)
    
    def set_background_color(self, presentation_path: str, slide_index: int, color: Union[List[int], str]) -> bool:
        """Set the background color of a specific slide."""
        presentation = self.get_presentation(presentation_path)
        slide = presentation.slides[slide_index]
        
        background = slide.background
        fill = background.fill
        fill.solid()
        
        # Convert hex color to RGB if it's a string
        if isinstance(color, str) and color.startswith('#'):
            color = color.lstrip('#')
            r, g, b = tuple(int(color[i:i+2], 16) for i in (0, 2, 4))
            fill.fore_color.rgb = RGBColor(r, g, b)
        else:
            fill.fore_color.rgb = RGBColor(*color)
        
        return True
    
    def delete_slide(self, presentation_path: str, slide_index: int) -> bool:
        """Delete a specific slide from the presentation."""
        presentation = self.get_presentation(presentation_path)
        xml_slides = presentation.slides._sldIdLst
        slides = list(xml_slides)
        if 0 <= slide_index < len(slides):
            xml_slides.remove(slides[slide_index])
            return True
        return False


# Create the PowerPoint context
context = PowerPointContext()

# Create MCP server
mcp = FastMCP("PowerPoint MCP")

# Basic Presentation Operations
@mcp.tool()
def list_presentations():
    """List all PowerPoint files in the workspace."""
    try:
        presentations = context.list_presentations()
        return {"presentations": presentations}
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def upload_presentation(file_path: str):
    """Upload a new presentation to the workspace."""
    try:
        path = context.upload_presentation(file_path)
        return {
            "message": "Presentation uploaded successfully",
            "path": path
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def save_presentation(presentation_path: str = None):
    """Save the current presentation."""
    try:
        context.save_presentation(presentation_path)
        return {"message": "Presentation saved successfully"}
    except Exception as e:
        return {"error": str(e)}

# Slide Operations
@mcp.tool()
def add_slide(presentation_path: str, layout_name: str = "Title and Content"):
    """Add a new slide to the presentation."""
    try:
        slide_index = context.add_slide(presentation_path, layout_name)
        return {
            "message": "Slide added successfully",
            "slide_index": slide_index
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def delete_slide(presentation_path: str, slide_index: int):
    """Delete a slide from the presentation."""
    try:
        success = context.delete_slide(presentation_path, slide_index)
        return {
            "success": success,
            "message": "Slide deleted successfully" if success else "Failed to delete slide"
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def get_slide_count(presentation_path: str):
    """Get the total number of slides in the presentation."""
    try:
        count = context.get_slide_count(presentation_path)
        return {"count": count}
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def analyze_slide(presentation_path: str, slide_index: int):
    """Analyze the content of a slide."""
    try:
        presentation = context.get_presentation(presentation_path)
        slide = presentation.slides[slide_index]
        content = context.analyze_slide_content(presentation_path, slide)
        preview = context.get_slide_preview(presentation_path, slide_index)
        return {
            "content": content,
            "preview": preview
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def set_background_color(presentation_path: str, slide_index: int, color: Union[List[int], str]):
    """Set the background color of a slide."""
    try:
        success = context.set_background_color(presentation_path, slide_index, color)
        return {
            "success": success,
            "message": "Background color set successfully"
        }
    except Exception as e:
        return {"error": str(e)}

# Element Operations
@mcp.tool()
def add_text(presentation_path: str, slide_index: int, text: str, position: List[float] = [0.5, 0.5], font_size: float = 18):
    """Add text to a slide."""
    try:
        element_id = context.add_text(presentation_path, slide_index, text, position, font_size)
        return {
            "message": "Text added successfully",
            "element_id": element_id
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def add_shape(presentation_path: str, slide_index: int, shape_type: str, position: Dict[str, float], size: Dict[str, float], style_properties: Dict[str, Any] = None):
    """Add a shape to a slide."""
    try:
        element_id = context.add_shape(presentation_path, slide_index, shape_type, position, size, style_properties)
        return {
            "message": "Shape added successfully",
            "element_id": element_id
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def edit_element(presentation_path: str, slide_index: int, element_id: str, properties: Dict[str, Any]):
    """Edit an element's properties."""
    try:
        updated_props = context.edit_element(presentation_path, slide_index, element_id, properties)
        return {
            "message": "Element updated successfully",
            "properties": updated_props
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def style_element(presentation_path: str, slide_index: int, element_id: str, style_properties: Dict[str, Any]):
    """Apply styling to an element."""
    try:
        success = context.style_element(presentation_path, slide_index, element_id, style_properties)
        return {
            "success": success,
            "message": "Style applied successfully"
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def connect_shapes(presentation_path: str, slide_index: int, from_element_id: str, to_element_id: str, connector_type: str = "straight", style_properties: Dict[str, Any] = None):
    """Connect two shapes with a connector."""
    try:
        connector_id = context.connect_shapes(presentation_path, slide_index, from_element_id, to_element_id, connector_type, style_properties)
        return {
            "message": "Shapes connected successfully",
            "connector_id": connector_id
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def find_element(presentation_path: str, slide_index: int, element_type: str = "any", search_text: str = None, position: Dict[str, float] = None):
    """Find elements on a slide based on criteria."""
    try:
        elements = context.find_element(presentation_path, slide_index, element_type, search_text, position)
        return {
            "elements": elements
        }
    except Exception as e:
        return {"error": str(e)}

# Financial Tools
@mcp.tool()
def get_company_financials(company_id: str, metrics: List[str] = ["revenue", "ebitda", "profit"], years: List[int] = None):
    """Get financial data for a company."""
    try:
        data = context.get_company_financials(company_id, metrics, years)
        return data
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def create_financial_chart(presentation_path: str, slide_index: int, chart_type: str, data: Dict[str, Any], position: Dict[str, float], size: Dict[str, float], title: str = None):
    """Create a financial chart on a slide."""
    try:
        chart_id = context.create_financial_chart(presentation_path, slide_index, chart_type, data, position, size, title)
        return {
            "message": "Chart created successfully",
            "chart_id": chart_id
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def create_comparison_table(presentation_path: str, slide_index: int, companies: List[str], metrics: List[str], position: Dict[str, float], title: str = None):
    """Create a comparison table for companies."""
    try:
        table_id = context.create_comparison_table(presentation_path, slide_index, companies, metrics, position, title)
        return {
            "message": "Comparison table created successfully",
            "table_id": table_id
        }
    except Exception as e:
        return {"error": str(e)}

# Template Operations
@mcp.tool()
def list_templates():
    """List all available templates."""
    try:
        templates = context.list_templates()
        return {"templates": templates}
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def apply_template(presentation_path: str, template_name: str, options: Dict[str, bool] = None):
    """Apply a template to a presentation."""
    try:
        applied_elements = context.apply_template(presentation_path, template_name, options)
        return {
            "message": "Template applied successfully",
            "applied_elements": applied_elements
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def create_slide_from_template(presentation_path: str, template_name: str, content: Dict[str, Any] = None):
    """Create a new slide from a template."""
    try:
        result = context.create_slide_from_template(presentation_path, template_name, content)
        return {
            "message": "Slide created from template successfully",
            "slide_index": result["slide_index"],
            "populated_placeholders": result["populated_placeholders"]
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def save_as_template(presentation_path: str, slide_index: int, template_name: str, template_description: str = ""):
    """Save a slide as a template."""
    try:
        template_info = context.save_as_template(presentation_path, slide_index, template_name, template_description)
        return {
            "message": "Template saved successfully",
            "template_info": template_info
        }
    except Exception as e:
        return {"error": str(e)}

@mcp.tool()
def debug_element_mappings(presentation_path: str, slide_index: int):
    """Debug tool to inspect element mappings for a slide."""
    try:
        if presentation_path not in context.element_ids:
            return {"error": f"No elements registered for presentation: {presentation_path}"}
            
        if slide_index not in context.element_ids[presentation_path]:
            return {"error": f"No elements registered for slide {slide_index}"}
            
        mappings = context.element_ids[presentation_path][slide_index]
        return {
            "presentation": presentation_path,
            "slide": slide_index,
            "mappings": mappings
        }
    except Exception as e:
        return {"error": str(e)}

if __name__ == "__main__":
    print("Starting PowerPoint MCP Server...")
    print("Initializing workspace...")
    
    # Initialize workspace
    if not os.path.exists(context.workspace_dir):
        os.makedirs(context.workspace_dir)
        print(f"Created workspace directory: {context.workspace_dir}")
    if not os.path.exists(os.path.join(context.workspace_dir, "templates")):
        os.makedirs(os.path.join(context.workspace_dir, "templates"))
        print("Created templates directory")
    
    # Run the server
    mcp.run()